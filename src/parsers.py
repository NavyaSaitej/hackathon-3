import gc
from pathlib import Path
from src.logger import logger
from src.exceptions import ParserFailureError


class IngestionRouter:
    """Routes files to the correct offline parser based on extension."""

    @staticmethod
    def process_file(filepath: str) -> str:
        path = Path(filepath)
        if not path.exists():
            logger.error(f"File not found: {filepath}")
            raise ParserFailureError(f"File not found: {filepath}")

        ext = path.suffix.lower()
        try:
            logger.info(f"Routing {filepath} to parser for extension '{ext}'")

            if ext == ".pdf":
                text = IngestionRouter._parse_pdf(path)
            elif ext in [".txt", ".md"]:
                text = path.read_text(encoding="utf-8")
            elif ext in [".wav", ".mp3"]:
                text = IngestionRouter._parse_audio(path)
            elif ext in [".png", ".jpg"]:
                text = IngestionRouter._parse_image(path)
            elif ext == ".docx":
                text = IngestionRouter._parse_docx(path)
            elif ext == ".pptx":
                text = IngestionRouter._parse_pptx(path)
            elif ext == ".xlsx":
                text = IngestionRouter._parse_xlsx(path)
            else:
                logger.warning(
                    f"Unsupported file type: {ext}. Attempting raw string extraction."
                )
                text = path.read_text(encoding="utf-8", errors="ignore")

            logger.info(f"Successfully extracted {len(text)} characters.")
            return text

        except Exception as e:
            logger.exception(f"Failed to parse {filepath}")
            raise ParserFailureError(f"Parsing failed for {filepath}") from e
        finally:
            logger.debug("Enforcing Garbage Collection after parser phase...")
            gc.collect()

    @staticmethod
    def _parse_pdf(path: Path) -> str:
        import fitz  # PyMuPDF

        logger.debug("Loading PyMuPDF...")
        doc = fitz.open(str(path))
        text = "".join([page.get_text() + "\n" for page in doc])
        return text

    @staticmethod
    def _parse_audio(path: Path) -> str:
        logger.debug("Loading faster-whisper...")
        try:
            from faster_whisper import WhisperModel
        except ImportError:
            logger.error("faster-whisper not installed.")
            raise ParserFailureError("Audio dependencies missing.")

        # Audio normalization via ffmpeg would go here in production
        model = WhisperModel("base", device="cpu", compute_type="int8")
        try:
            segments, _ = model.transcribe(str(path), beam_size=5)
            text = " ".join([segment.text for segment in segments])
        finally:
            logger.debug("Purging faster-whisper from memory.")
            del model
        return text

    @staticmethod
    def _parse_image(path: Path) -> str:
        logger.debug("Loading rapidocr_onnxruntime...")
        try:
            from rapidocr_onnxruntime import RapidOCR
        except ImportError:
            logger.error("rapidocr_onnxruntime not installed.")
            raise ParserFailureError("Image dependencies missing.")

        engine = RapidOCR()
        try:
            result, _ = engine(str(path))
            if result:
                text = " ".join([res[1] for res in result])
            else:
                text = ""
        finally:
            logger.debug("Purging RapidOCR from memory.")
            del engine
        return text

    @staticmethod
    def _parse_docx(path: Path) -> str:
        logger.debug("Loading python-docx...")
        try:
            from docx import Document
        except ImportError:
            raise ParserFailureError("python-docx not installed.")
        doc = Document(str(path))
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _parse_pptx(path: Path) -> str:
        logger.debug("Loading python-pptx...")
        try:
            from pptx import Presentation
        except ImportError:
            raise ParserFailureError("python-pptx not installed.")
        prs = Presentation(str(path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def _parse_xlsx(path: Path) -> str:
        logger.debug("Loading pandas for excel extraction...")
        try:
            import pandas as pd
        except ImportError:
            raise ParserFailureError("pandas not installed.")
        df = pd.read_excel(str(path))
        return df.to_string()
