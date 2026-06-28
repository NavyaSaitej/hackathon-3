import pandas as pd
from docx import Document
from pptx import Presentation
import fitz
from pathlib import Path

def generate():
    Path("dummy_files").mkdir(exist_ok=True)
    
    # Text and Markdown
    Path("dummy_files/test.txt").write_text("Hello from txt", encoding="utf-8")
    Path("dummy_files/test.md").write_text("# Hello from md", encoding="utf-8")
    
    # Docx
    doc = Document()
    doc.add_paragraph("Hello from docx")
    doc.save("dummy_files/test.docx")
    
    # PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello from pptx"
    prs.save("dummy_files/test.pptx")
    
    # XLSX
    df = pd.DataFrame({"Message": ["Hello from xlsx"]})
    df.to_excel("dummy_files/test.xlsx", index=False)
    
    # PDF
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((50, 50), "Hello from pdf")
    pdf.save("dummy_files/test.pdf")
    pdf.close()
    
    # Dummy image/audio
    Path("dummy_files/test.wav").write_bytes(b"dummy audio")
    Path("dummy_files/test.png").write_bytes(b"dummy image")
    
    print("Files generated!")

if __name__ == "__main__":
    generate()
